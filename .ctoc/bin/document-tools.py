#!/usr/bin/env python3
"""
CTOC Document Tools - Seamless document creation and reading

This script handles dependency installation and provides a unified interface
for document operations. Users don't need to know about Python or pip.

Usage:
    python document-tools.py check          # Check if dependencies are installed
    python document-tools.py install        # Install all document dependencies
    python document-tools.py create <type> <output> [options]
    python document-tools.py read <file>

Examples:
    python document-tools.py check
    python document-tools.py install
    python document-tools.py read report.pdf
"""

import sys
import subprocess
import json
from pathlib import Path

# Required packages for each feature
PACKAGES = {
    'pdf_write': ['fpdf2', 'reportlab', 'Pillow'],
    'pdf_read': ['PyPDF2', 'pdfplumber'],
    'docx': ['python-docx', 'Pillow'],
    'pptx': ['python-pptx', 'Pillow'],
    'charts': ['matplotlib', 'pandas'],
}

ALL_PACKAGES = list(set(pkg for pkgs in PACKAGES.values() for pkg in pkgs))


def check_python():
    """Check if Python 3 is available."""
    try:
        version = sys.version_info
        if version.major >= 3 and version.minor >= 8:
            return True, f"Python {version.major}.{version.minor}.{version.micro}"
        return False, f"Python {version.major}.{version.minor} (need 3.8+)"
    except Exception as e:
        return False, str(e)


def check_package(package):
    """Check if a package is installed."""
    try:
        # Map package names to import names
        import_map = {
            'python-docx': 'docx',
            'python-pptx': 'pptx',
            'fpdf2': 'fpdf',
            'PyPDF2': 'PyPDF2',
            'Pillow': 'PIL',
        }
        import_name = import_map.get(package, package)
        __import__(import_name)
        return True
    except ImportError:
        return False


def install_packages(packages):
    """Install packages using pip."""
    print(f"Installing document tools dependencies...")

    # Try different pip commands
    pip_commands = [
        [sys.executable, '-m', 'pip', 'install', '--user', '--quiet'],
        ['pip3', 'install', '--user', '--quiet'],
        ['pip', 'install', '--user', '--quiet'],
    ]

    for pip_cmd in pip_commands:
        try:
            cmd = pip_cmd + packages
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode == 0:
                print(f"  Installed: {', '.join(packages)}")
                return True
        except FileNotFoundError:
            continue

    print(f"  Failed to install packages. Please run:")
    print(f"    pip install --user {' '.join(packages)}")
    return False


def check_dependencies():
    """Check all dependencies and return status."""
    status = {
        'python': check_python(),
        'packages': {},
        'ready': True
    }

    for pkg in ALL_PACKAGES:
        installed = check_package(pkg)
        status['packages'][pkg] = installed
        if not installed:
            status['ready'] = False

    return status


def ensure_dependencies(feature=None):
    """Ensure dependencies are installed for a feature."""
    if feature:
        packages = PACKAGES.get(feature, [])
    else:
        packages = ALL_PACKAGES

    missing = [pkg for pkg in packages if not check_package(pkg)]

    if missing:
        return install_packages(missing)
    return True


def cmd_check():
    """Check command - show dependency status."""
    status = check_dependencies()

    py_ok, py_version = status['python']
    print(f"\nCTOC Document Tools - Dependency Check")
    print(f"{'=' * 40}")
    print(f"Python: {'OK' if py_ok else 'MISSING'} ({py_version})")
    print()
    print("Packages:")

    for pkg, installed in status['packages'].items():
        icon = '' if installed else ''
        print(f"  {icon} {pkg}")

    print()
    if status['ready']:
        print("All dependencies installed. Ready to create documents!")
    else:
        print("Missing dependencies. Run: python document-tools.py install")

    return 0 if status['ready'] else 1


def cmd_install():
    """Install command - install all dependencies."""
    print("\nCTOC Document Tools - Installing Dependencies")
    print(f"{'=' * 40}")

    if ensure_dependencies():
        print("\nAll document tools ready!")
        return 0
    else:
        print("\nSome packages failed to install.")
        return 1


def cmd_read(filepath):
    """Read command - extract content from a document."""
    path = Path(filepath)

    if not path.exists():
        print(f"Error: File not found: {filepath}")
        return 1

    ext = path.suffix.lower()

    if ext == '.pdf':
        if not ensure_dependencies('pdf_read'):
            return 1
        return read_pdf(path)
    elif ext == '.docx':
        if not ensure_dependencies('docx'):
            return 1
        return read_docx(path)
    elif ext == '.pptx':
        if not ensure_dependencies('pptx'):
            return 1
        return read_pptx(path)
    else:
        print(f"Error: Unsupported file type: {ext}")
        print("Supported: .pdf, .docx, .pptx")
        return 1


def read_pdf(path):
    """Read PDF and output content."""
    try:
        import pdfplumber

        result = {
            'file': str(path),
            'type': 'pdf',
            'pages': [],
            'tables': []
        }

        with pdfplumber.open(path) as pdf:
            result['page_count'] = len(pdf.pages)

            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ''
                tables = page.extract_tables() or []

                result['pages'].append({
                    'number': i,
                    'text': text[:2000] + ('...' if len(text) > 2000 else '')
                })

                for table in tables:
                    result['tables'].append({
                        'page': i,
                        'rows': len(table),
                        'data': table[:5]  # First 5 rows
                    })

        print(json.dumps(result, indent=2))
        return 0

    except Exception as e:
        print(f"Error reading PDF: {e}")
        return 1


def read_docx(path):
    """Read DOCX and output content."""
    try:
        from docx import Document

        doc = Document(path)

        result = {
            'file': str(path),
            'type': 'docx',
            'paragraphs': [],
            'headings': [],
            'tables': []
        }

        for para in doc.paragraphs:
            if para.style.name.startswith('Heading'):
                level = para.style.name[-1] if para.style.name[-1].isdigit() else '0'
                result['headings'].append({
                    'level': int(level),
                    'text': para.text
                })
            elif para.text.strip():
                result['paragraphs'].append(para.text[:500])

        for table in doc.tables:
            rows = []
            for row in table.rows[:10]:  # First 10 rows
                rows.append([cell.text for cell in row.cells])
            result['tables'].append({'rows': rows})

        print(json.dumps(result, indent=2))
        return 0

    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return 1


def read_pptx(path):
    """Read PPTX and output content."""
    try:
        from pptx import Presentation

        prs = Presentation(path)

        result = {
            'file': str(path),
            'type': 'pptx',
            'slide_count': len(prs.slides),
            'slides': []
        }

        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                'number': i,
                'title': '',
                'content': [],
                'notes': ''
            }

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title
                                slide_data['title'] = text
                            else:
                                slide_data['content'].append(text)
                        else:
                            slide_data['content'].append(text)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                slide_data['notes'] = notes[:500] if notes else ''

            result['slides'].append(slide_data)

        print(json.dumps(result, indent=2))
        return 0

    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return 1


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        return 0

    command = sys.argv[1]

    if command == 'check':
        return cmd_check()
    elif command == 'install':
        return cmd_install()
    elif command == 'read':
        if len(sys.argv) < 3:
            print("Usage: document-tools.py read <filepath>")
            return 1
        return cmd_read(sys.argv[2])
    else:
        print(f"Unknown command: {command}")
        print("Commands: check, install, read")
        return 1


if __name__ == '__main__':
    sys.exit(main())
